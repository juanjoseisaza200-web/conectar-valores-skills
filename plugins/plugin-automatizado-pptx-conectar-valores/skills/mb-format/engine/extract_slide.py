#!/usr/bin/env python3
"""
extract_slide.py — Lee UNA slide de un .pptx y devuelve JSON estructurado.

Uso: python3 extract_slide.py <archivo.pptx> <slide_idx_0based>

Salida JSON:
{
  "slideW": float,   // ancho en inches
  "slideH": float,   // alto en inches
  "fragRatio": float, // ratio de fragmentacion de texto (señal de calidad)
  "elements": [
    {
      "role": str,   // clasificacion preliminar del elemento
      "text": str,   // texto plano concatenado (si aplica)
      "x": float, "y": float, "w": float, "h": float,
      "meta": {}     // datos extra según el tipo (filas de tabla, n_kpis, etc.)
    }
  ]
}

Roles que puede asignar:
  title           — elemento más grande, posicionado arriba
  subtitle        — debajo del título, fuente menor, color distinto
  company         — nombre de empresa (en este template específico, viene subrayado)
  description     — bloque de texto descriptivo, izquierda
  fin_table       — tabla financiera (estado de resultados / balance)
  kpi_group       — grupo de tarjetas KPI pequeñas (detectado como grupo de rectangulos+texto)
  chart           — chart nativo de PowerPoint (no lo tocamos, solo lo registramos)
  closing_text    — texto de cierre/conclusión (bloque de texto bajo el chart)
  section_header  — título de sección dentro de la slide ("Valoración", "Información financiera", etc.)
  footer_ignore   — elementos de footer que mb-format reemplaza completamente
  unknown         — no clasificado, requiere revisión manual
"""

import sys
import json
from pptx import Presentation

EMU = 914400  # 1 inch en EMU

def safe_text(shape):
    try:
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
    except Exception:
        pass
    return ""

def safe_font_size(shape):
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        return run.font.size.pt
    except Exception:
        pass
    return None

def safe_bold(shape):
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    return bool(run.font.bold)
    except Exception:
        pass
    return False

def count_runs(shape):
    try:
        if shape.has_text_frame:
            return sum(len(p.runs) for p in shape.text_frame.paragraphs)
    except Exception:
        pass
    return 0

def word_count(text):
    return len(text.split()) if text else 0

def extract_table(shape):
    try:
        tbl = shape.table
        rows = []
        for row in tbl.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        return rows
    except Exception:
        return []

def walk_shape(shape, elements, slide_h_in, depth=0):
    """Recorre shapes recursivamente, expandiendo GROUPs."""
    if shape.shape_type == 6:  # GROUP
        # Un grupo puede ser un conjunto de KPIs — lo detectamos por
        # sus hijos (múltiples rectángulos + textos pequeños)
        try:
            child_types = [s.shape_type for s in shape.shapes]
            has_rect = child_types.count(1) >= 2  # MSO_SHAPE_TYPE.AUTO_SHAPE = 1
            child_texts = [safe_text(s) for s in shape.shapes if safe_text(s)]
            # Heurística KPI: varios rectángulos con textos cortos (valores numéricos)
            short_texts = [t for t in child_texts if len(t) < 30]
            if has_rect and len(short_texts) >= 3 and depth <= 1:
                x = shape.left / EMU if shape.left else 0
                y = shape.top / EMU if shape.top else 0
                w = shape.width / EMU if shape.width else 0
                h = shape.height / EMU if shape.height else 0
                elements.append({
                    "role": "kpi_group",
                    "text": " | ".join(short_texts),
                    "x": x, "y": y, "w": w, "h": h,
                    "meta": { "kpi_texts": child_texts, "n_children": len(child_types) }
                })
                return  # no expandir más adentro
        except Exception:
            pass
        # Si no es un grupo de KPIs, expandir recursivamente
        try:
            for sub in shape.shapes:
                walk_shape(sub, elements, slide_h_in, depth+1)
        except Exception:
            pass
        return

    x = (shape.left or 0) / EMU
    y = (shape.top or 0) / EMU
    w = (shape.width or 0) / EMU
    h = (shape.height or 0) / EMU
    text = safe_text(shape)
    font_size = safe_font_size(shape)
    is_bold = safe_bold(shape)

    # ── CHART nativo ──────────────────────────────────────────────────────────
    if shape.shape_type == 3:  # CHART
        elements.append({
            "role": "chart",
            "text": "",
            "x": x, "y": y, "w": w, "h": h,
            "meta": { "shape_name": shape.name }
        })
        return

    # ── TABLA ─────────────────────────────────────────────────────────────────
    if shape.has_table:
        rows = extract_table(shape)
        elements.append({
            "role": "fin_table",
            "text": " | ".join([" ".join(r) for r in rows[:2]]),  # preview
            "x": x, "y": y, "w": w, "h": h,
            "meta": { "rows": rows }
        })
        return

    # ── PICTURE / LOGO ────────────────────────────────────────────────────────
    if shape.shape_type == 13:  # PICTURE
        # Logo en esquina superior derecha → reemplazar por CV
        # Logo en esquina inferior derecha → mantener (ya es CV)
        is_footer_zone = y > slide_h_in * 0.8
        elements.append({
            "role": "footer_ignore" if is_footer_zone else "logo_replace",
            "text": "",
            "x": x, "y": y, "w": w, "h": h,
            "meta": { "shape_name": shape.name }
        })
        return

    # ── LÍNEA DECORATIVA ─────────────────────────────────────────────────────
    if shape.shape_type == 9:  # LINE
        is_footer_zone = y > slide_h_in * 0.85
        elements.append({
            "role": "footer_ignore" if is_footer_zone else "decorative_line",
            "text": "",
            "x": x, "y": y, "w": w, "h": h,
            "meta": {}
        })
        return

    # ── TEXTO ─────────────────────────────────────────────────────────────────
    if not text:
        return  # shape sin texto visible — ignorar

    # footer conocidos
    footer_phrases = ["privado y confidencial", "confidencial", "pagina", "página"]
    if any(fp in text.lower() for fp in footer_phrases) and y > slide_h_in * 0.8:
        elements.append({
            "role": "footer_ignore",
            "text": text, "x": x, "y": y, "w": w, "h": h, "meta": {}
        })
        return

    # Clasificación por posición y características de texto
    is_top = y < slide_h_in * 0.25
    is_large_font = font_size is not None and font_size >= 20
    words = word_count(text)
    is_short = words <= 8
    is_underlined = False
    try:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.underline:
                        is_underlined = True
    except Exception:
        pass

    # secciones internas tipo "Valoración", "Información financiera"
    section_keywords = ["valoración", "valoracion", "información financiera",
                        "informacion financiera", "aspectos relevantes"]
    if any(kw in text.lower() for kw in section_keywords) and is_short and is_bold:
        elements.append({
            "role": "section_header",
            "text": text, "x": x, "y": y, "w": w, "h": h, "meta": {}
        })
        return

    # título principal (grande, arriba, corto)
    if is_top and is_large_font and is_short:
        elements.append({
            "role": "title",
            "text": text, "x": x, "y": y, "w": w, "h": h,
            "meta": { "font_size": font_size }
        })
        return

    # subtítulo (arriba, fuente media, debajo del título)
    if is_top and not is_large_font and is_short and not is_underlined:
        elements.append({
            "role": "subtitle",
            "text": text, "x": x, "y": y, "w": w, "h": h,
            "meta": { "font_size": font_size }
        })
        return

    # nombre de empresa (corto, subrayado, posición media-alta)
    if is_underlined and is_short and y < slide_h_in * 0.35:
        elements.append({
            "role": "company",
            "text": text, "x": x, "y": y, "w": w, "h": h, "meta": {}
        })
        return

    # texto de cierre / conclusión (bloques largos en la parte baja de la columna izquierda)
    if words > 30 and y > slide_h_in * 0.6:
        elements.append({
            "role": "closing_text",
            "text": text, "x": x, "y": y, "w": w, "h": h, "meta": {}
        })
        return

    # descripción general (bloque de texto largo, columna izquierda)
    if words > 20:
        elements.append({
            "role": "description",
            "text": text, "x": x, "y": y, "w": w, "h": h, "meta": {}
        })
        return

    # todo lo demás sin clasificar
    elements.append({
        "role": "unknown",
        "text": text, "x": x, "y": y, "w": w, "h": h,
        "meta": { "font_size": font_size, "bold": is_bold, "words": words }
    })


def extract(pptx_path, slide_idx):
    p = Presentation(pptx_path)
    slide_w = p.slide_width / EMU
    slide_h = p.slide_height / EMU
    slide = p.slides[slide_idx]

    elements = []
    total_runs = 0
    total_words = 0

    for shape in slide.shapes:
        walk_shape(shape, elements, slide_h)

    # calcular ratio de fragmentación para reporting
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                txt = shape.text_frame.text
                total_runs += sum(len(p.runs) for p in shape.text_frame.paragraphs)
                total_words += word_count(txt)
        except Exception:
            pass

    frag_ratio = total_runs / max(total_words, 1)

    # ordenar por posición Y para facilitar lectura del output
    elements.sort(key=lambda e: e["y"])

    return {
        "slideW": slide_w,
        "slideH": slide_h,
        "fragRatio": round(frag_ratio, 2),
        "elements": elements
    }


if __name__ == "__main__":
    pptx_path = sys.argv[1]
    slide_idx = int(sys.argv[2])
    result = extract(pptx_path, slide_idx)
    print(json.dumps(result, ensure_ascii=False, indent=2))
