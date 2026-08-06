"""Deck style — helpers para construir slides en estilo MUFG-elegant.
Arial typography, navy+teal sage palette, generous margins, minimalist tables.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Palette (extracted from MUFG reference)
NAVY        = RGBColor(0x1F, 0x3D, 0x5C)   # primary titles, table headers
NAVY_DARK   = RGBColor(0x14, 0x2A, 0x40)
TEAL        = RGBColor(0x5D, 0x9E, 0xA7)   # accent — subtle use
TEAL_DARK   = RGBColor(0x3E, 0x7A, 0x82)
SAGE        = RGBColor(0xA8, 0xC5, 0xC0)   # very light teal-sage
SAGE_BG     = RGBColor(0xE6, 0xEE, 0xEC)   # background sage tint
GRAY_LIGHT  = RGBColor(0xF4, 0xF6, 0xF7)   # alt row, sidebar bg
GRAY_LINE   = RGBColor(0xD8, 0xDC, 0xDF)
GRAY_MID    = RGBColor(0xA9, 0xAE, 0xB5)
GRAY_TEXT   = RGBColor(0x6B, 0x71, 0x79)
DARK_TEXT   = RGBColor(0x2A, 0x2D, 0x33)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CREAM       = RGBColor(0xF6, 0xE8, 0xCB)   # highlight rows (totals)
ORANGE      = RGBColor(0xE0, 0x8E, 0x3C)
GREEN       = RGBColor(0x6F, 0xA0, 0x6F)
GOLD        = RGBColor(0xC9, 0xA2, 0x27)

# Font
FONT = "Arial"

# Nombre del proyecto en el footer. Sobrescribir por proyecto:
#   import deck_style; deck_style.PROJECT_NAME = "NOMBRE DEL ACTIVO"
PROJECT_NAME = "PROYECTO"


# Margins (inches)
LM = 0.65   # left margin
RM = 0.65   # right margin
TM = 1.15   # top margin (below master title bar)
BM = 0.40   # bottom margin (above footer)
SLIDE_W = 13.333
SLIDE_H = 7.5
USABLE_W = SLIDE_W - LM - RM
USABLE_H = SLIDE_H - TM - BM


# ---------- Helpers ----------
def set_subtitle(slide, text):
    """Find Rectangle 90 banner shape by name and replace subtitle 'xxxx'."""
    rect = None
    for sh in slide.shapes:
        if sh.name == "Rectangle 90" and sh.has_text_frame:
            tf = sh.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                txt = tf.paragraphs[0].runs[0].text or ""
                if "xxxx" in txt or txt.strip() == "" or rect is None:
                    rect = sh
                    if "xxxx" in txt:
                        break
    if rect is None:
        return
    p = rect.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.add_run().text = text


def add_textbox(slide, left, top, width, height, text, *,
                size=10, bold=False, italic=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, vanchor=MSO_ANCHOR.TOP,
                font=FONT, line_spacing=1.20):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = vanchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_rich(slide, left, top, width, height, paragraphs, *,
             vanchor=MSO_ANCHOR.TOP, line_spacing=1.20):
    """paragraphs = list of lists of (text, props_dict). Each inner list = one paragraph."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = vanchor
    for pi, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        if runs and "align" in runs[0]:
            p.alignment = runs[0]["align"]
        space_before = None
        for txt, props in runs:
            if "space_before" in props:
                space_before = props["space_before"]
            r = p.add_run()
            r.text = txt
            r.font.name = props.get("font", FONT)
            r.font.size = Pt(props.get("size", 10))
            r.font.bold = props.get("bold", False)
            r.font.italic = props.get("italic", False)
            r.font.color.rgb = props.get("color", DARK_TEXT)
        if space_before is not None:
            p.space_before = Pt(space_before)
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=None, line_w=0.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(height))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_filled_text(slide, left, top, width, height, text, *,
                    fill=NAVY, color=WHITE, size=10, bold=True,
                    align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE,
                    font=FONT, line=None, line_spacing=1.10,
                    margin=0.05):
    sh = add_rect(slide, left, top, width, height, fill=fill, line=line)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = vanchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line_txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line_txt
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return sh


def add_section_title(slide, left, top, width, title, *,
                       title_size=11, color=NAVY, with_line=True):
    """Section header in MUFG style: small text + thin teal line below."""
    add_textbox(slide, left, top, width, 0.28, title.upper(),
                size=title_size, bold=True, color=color,
                font=FONT)
    if with_line:
        # Thin teal accent line under
        add_rect(slide, left, top + 0.30, 0.50, 0.025,
                 fill=TEAL, line=None)


def add_horizontal_divider(slide, left, top, width, color=GRAY_LINE, weight=0.75):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(left), Inches(top),
                                       Inches(left + width), Inches(top))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_table_mufg(slide, left, top, width, height, data, *,
                   header_fill=NAVY, header_color=WHITE,
                   row_alt=GRAY_LIGHT, total_rows=None,
                   highlight_rows=None,
                   header_size=10, body_size=9.5,
                   col_aligns=None,
                   first_col_bold=False):
    rows = len(data); cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols,
                                        Inches(left), Inches(top),
                                        Inches(width), Inches(height))
    tbl = tbl_shape.table
    total_rows = total_rows or []
    highlight_rows = highlight_rows or []
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            if col_aligns:
                p.alignment = col_aligns[c]
            else:
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            # Fill
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            elif r in total_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = CREAM
            elif r in highlight_rows:
                cell.fill.solid(); cell.fill.fore_color.rgb = SAGE_BG
            else:
                if r % 2 == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = GRAY_LIGHT
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.name = FONT
            run.font.size = Pt(header_size if r == 0 else body_size)
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = header_color
            elif r in total_rows:
                run.font.bold = True
                run.font.color.rgb = NAVY
            else:
                run.font.bold = (c == 0 and first_col_bold)
                run.font.color.rgb = DARK_TEXT
    # Remove inner vertical borders by setting cells border properties via XML
    # python-pptx doesn't expose this cleanly; default tbl style may include borders
    # Will rely on visual contrast (no manual XML manipulation here for safety)
    return tbl_shape


def add_kpi_strip(slide, left, top, width, height, kpis, *,
                  fill=TEAL, color=WHITE, num_size=22, lbl_size=9):
    """Bottom KPI strip — MUFG style. kpis = list of (value, label)."""
    add_rect(slide, left, top, width, height, fill=fill, line=None)
    n = len(kpis)
    cell_w = width / n
    for i, (val, lbl) in enumerate(kpis):
        x = left + i * cell_w
        # Big number
        add_textbox(slide, x, top + 0.10, cell_w, height * 0.55,
                    val, size=num_size, bold=True, color=color,
                    align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)
        # Label
        add_textbox(slide, x, top + height * 0.62, cell_w, height * 0.32,
                    lbl, size=lbl_size, bold=False, color=color,
                    align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)
        # Vertical separators between KPIs
        if i > 0:
            sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                              Inches(x), Inches(top + 0.20),
                                              Inches(x), Inches(top + height - 0.20))
            sep.line.color.rgb = WHITE
            sep.line.width = Pt(0.5)


def add_kpi_sidebar(slide, left, top, width, height, kpis, *,
                    fill=GRAY_LIGHT, num_color=NAVY,
                    num_size=24, lbl_size=9):
    """Right-side stacked KPI list (MUFG p10 style). kpis = list of (val, label)."""
    add_rect(slide, left, top, width, height, fill=fill, line=None)
    n = len(kpis)
    item_h = height / n
    for i, (val, lbl) in enumerate(kpis):
        y = top + i * item_h
        add_textbox(slide, left + 0.20, y + 0.15, width - 0.40, item_h * 0.45,
                    val, size=num_size, bold=True, color=num_color,
                    align=PP_ALIGN.LEFT, vanchor=MSO_ANCHOR.TOP)
        add_textbox(slide, left + 0.20, y + 0.15 + item_h * 0.45,
                    width - 0.40, item_h * 0.40,
                    lbl, size=lbl_size, color=GRAY_TEXT,
                    align=PP_ALIGN.LEFT, vanchor=MSO_ANCHOR.TOP,
                    line_spacing=1.15)


def add_footer(slide, page_num, total_pages=None, project=None):
    """MUFG-style footer at bottom. Subtle, all caps wide-tracking."""
    project = project or PROJECT_NAME
    fy = 7.20
    add_textbox(slide, 0.40, fy, 4.0, 0.20,
                "ESTRICTAMENTE PRIVADO Y CONFIDENCIAL",
                size=8, bold=True, color=NAVY, font=FONT)
    add_textbox(slide, 4.50, fy, 4.30, 0.20,
                "PREPARADO POR  ·  CONECTAR VALORES S.A.S.",
                size=8, color=GRAY_TEXT, align=PP_ALIGN.CENTER, font=FONT)
    add_textbox(slide, 9.0, fy, 3.90, 0.20,
                f"{project}  ·  {page_num:02d}",
                size=8, bold=True, color=NAVY, align=PP_ALIGN.RIGHT, font=FONT)


def add_section_chip(slide, left, top, text, color=TEAL):
    """Small ALL-CAPS section chip with letter-spacing — MUFG style."""
    add_textbox(slide, left, top, 5.0, 0.22, text.upper(),
                size=8, bold=True, color=color, font=FONT)
    add_rect(slide, left, top + 0.24, 0.40, 0.02, fill=color, line=None)


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.2, head=True):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1),
                                       Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if head:
        ln = conn.line._get_or_add_ln()
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('h', 'med')
    return conn
